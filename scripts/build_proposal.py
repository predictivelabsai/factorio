"""Build the Universalbank proposal in four formats, per language (en / ru).

Author: Consistente Ltd. Platform: Factorio. Audience: Universalbank.

Writes, into docs/, for each language:
    universalbank_consistente_proposal_<lang>.md     (source deck; mermaid diagrams)
    universalbank_consistente_proposal_<lang>.html   (markdown + assets/proposal.css; PNG diagrams)
    universalbank_consistente_proposal_<lang>.pdf     (WeasyPrint, A4 landscape slides)
    universalbank_consistente_proposal_<lang>.pptx    (python-pptx, 16:9 deck)

Architecture diagrams are authored once as Mermaid and rendered to PNG via the
kroki.io service (so the PDF/PPTX show real diagrams and the .md keeps the
mermaid source, which renders natively on GitHub). Content is authored once,
bilingually, so all formats and both languages stay in sync.

    python -m scripts.build_proposal
"""

from __future__ import annotations

import urllib.request
from pathlib import Path

import markdown as md_lib
from weasyprint import HTML
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from scripts.diagrams import BILINGUAL, render_bilingual

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
IMG = DOCS / "img"
DIAG = IMG / "proposal"
CSS = DOCS / "assets" / "proposal.css"

# Brand
GREEN = RGBColor(0x1F, 0x5D, 0x43)
INK = RGBColor(0x14, 0x23, 0x1B)
MUTED = RGBColor(0x41, 0x50, 0x46)
GOLD = RGBColor(0xC8, 0xA2, 0x4B)
PARCH = RGBColor(0xF7, 0xF6, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Architecture diagrams (Mermaid, rendered to PNG via kroki) ──────────────

_INIT = "%%{init: {'theme':'neutral'}}%%\n"

DIAGRAMS = {
    "d1-architecture": _INIT + """flowchart TB
    USER(["Sellers & Investors<br/>(browser · trilingual)"])
    subgraph APP["Factorio — one FastHTML process"]
        LAND["Landing site"]
        PROD["Investor app<br/>dashboard · marketplace · portfolio"]
        AITRI["AI triage<br/>/app/triage"]
        AIREP["AI reporting<br/>/app/assistant"]
        SCORE["Credit-scoring engine"]
        SPVM["SPV / investor module"]
    end
    subgraph DATA["PostgreSQL — factorio schema"]
        DB[("invoices · funding · investments<br/>companies · settlements<br/>ai_* · credit_*")]
    end
    subgraph EXT["External services"]
        XAI[["Grok LLM<br/>x.ai"]]
        SOLIQ[["SoliqOnline / Didox<br/>e-invoice API"]]
        OB[["Open banking<br/>bank transaction data"]]
        CRIF[["CRIF credit bureau"]]
        CBU[["CBU collateral registry"]]
        CUST[["DIFC SPV · custodian<br/>international investors"]]
    end
    USER --> LAND
    USER --> PROD
    USER --> AITRI
    USER --> AIREP
    AITRI --> XAI
    AIREP --> XAI
    PROD --> DB
    AITRI --> DB
    AIREP --> DB
    SCORE --> SOLIQ
    SCORE --> OB
    SCORE --> CRIF
    SCORE --> DB
    PROD --> CBU
    SPVM --> CUST
    SPVM --> DB
    style XAI fill:#1F5D43,color:#fff
    style SOLIQ fill:#2563EB,color:#fff
    style OB fill:#2563EB,color:#fff
    style CRIF fill:#2563EB,color:#fff
    style CUST fill:#C8A24B,color:#14231B
    style AITRI fill:#DCF3E8
    style AIREP fill:#DCF3E8
    style SCORE fill:#DCF3E8
""",
    "d2-ai-chat": _INIT + """flowchart LR
    S["Seller message<br/>(natural language)"] --> TRI["Triage assistant<br/>app_routes/assistant.py"]
    I["Investor question"] --> REP["Reporting assistant"]
    TRI --> TOOLS["Tools<br/>lookup_debtor · get_soliq_invoice<br/>indicative_terms"]
    REP --> CTX["Grounding<br/>investor's live positions"]
    TOOLS --> GROK[["Grok · x.ai"]]
    CTX --> GROK
    GROK --> OUT1["Indicative grade A–D<br/>advance rate · next documents"]
    GROK --> OUT2["Grounded portfolio answer<br/>(figures trace to real data)"]
    style GROK fill:#1F5D43,color:#fff
    style OUT1 fill:#DCF3E8
    style OUT2 fill:#DCF3E8
""",
    "d3-scoring": _INIT + """flowchart LR
    A["SoliqOnline<br/>invoice + tax history"] --> FE["Feature engineering"]
    B["Open banking<br/>cash-flow signals"] --> FE
    C["CRIF bureau<br/>obligations · defaults"] --> FE
    FE --> M["Scoring model<br/>gradient-boosted + rules"]
    M --> D["Risk grade A–D<br/>advance rate · price"]
    D --> L["Grok<br/>plain-language rationale<br/>+ adverse-action reasons"]
    D --> AUD[("Decision audit<br/>versioned · reproducible")]
    style M fill:#1F5D43,color:#fff
    style D fill:#C8A24B,color:#14231B
    style L fill:#DCF3E8
""",
    "d4-spv": _INIT + """flowchart TB
    POOL[("Uzbek invoice pool<br/>SoliqOnline-verified · buyer-confirmed")]
    FACT["Factorio platform<br/>Universalbank — originator & servicer"]
    SPV["DIFC SPV<br/>(bankruptcy-remote)"]
    subgraph INV["International investors"]
        CONV["Conventional<br/>note / participation"]
        ISL["Shariah tracks<br/>Murabaha → Wakala → Sukuk"]
    end
    POOL --> FACT --> SPV
    SPV --> CONV
    SPV --> ISL
    CONV -->|USD capital| SPV
    ISL -->|USD capital| SPV
    SPV -->|funding| FACT
    FACT -->|collections| SPV
    SPV -->|profit distribution| INV
    style SPV fill:#1F5D43,color:#fff
    style ISL fill:#C8A24B,color:#14231B
    style POOL fill:#2563EB,color:#fff
""",
    "d5-roadmap": _INIT + """flowchart LR
    P0["Phase 0 · now<br/>Working AI prototype<br/>triage + reporting shipped"] --> P1["Phase 1 · M1–3<br/>Platform GA + streaming<br/>SoliqOnline integration"]
    P1 --> P2["Phase 2 · M3–6<br/>Open-banking credit scoring<br/>volume pricing live"]
    P2 --> P3["Phase 3 · M6–12<br/>DIFC SPV · first int'l capital<br/>Murabaha pilot"]
    P3 --> P4["Phase 4 · Y2+<br/>Wakala fund → Sukuk<br/>regional scale"]
    style P0 fill:#DCF3E8
    style P2 fill:#1F5D43,color:#fff
    style P4 fill:#C8A24B,color:#14231B
""",
}


def render_diagrams() -> None:
    """Render each Mermaid diagram to PNG via kroki.io. Best-effort: a failed
    render prints a warning and leaves that slide image missing rather than
    aborting the whole build."""
    DIAG.mkdir(parents=True, exist_ok=True)
    for key, src in DIAGRAMS.items():
        out = DIAG / f"{key}.png"
        try:
            req = urllib.request.Request(
                "https://kroki.io/mermaid/png",
                data=src.encode("utf-8"),
                headers={"Content-Type": "text/plain",
                         "User-Agent": "Mozilla/5.0 (Factorio proposal builder)"},
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=60) as resp:
                out.write_bytes(resp.read())
            print(f"  diagram {out.relative_to(ROOT)}  ({out.stat().st_size/1024:.0f} KB)")
        except Exception as e:  # noqa: BLE001 - best effort
            print(f"  WARNING: could not render {key}: {e}")


# ── Content (authored once, bilingual) ──────────────────────────────────────

TITLE = {
    "en": "Factorio for Universalbank",
    "ru": "Factorio для Universalbank",
}
SUBTITLE = {
    "en": "AI-native invoice financing · an international-investor SPV · UAE entry",
    "ru": "Факторинг на базе ИИ · SPV для международных инвесторов · выход в ОАЭ",
}
INTRO = {
    "en": ("A proposal by **Consistente Ltd** to build and operate a next-generation, "
           "white-label invoice-financing platform for Universalbank — more efficient "
           "than the incumbent, priced on volume alone, with an AI core, an "
           "international-investor SPV, and a route into the UAE."),
    "ru": ("Предложение от **Consistente Ltd** построить и эксплуатировать для "
           "Universalbank платформу факторинга нового поколения (white-label) — "
           "эффективнее действующего решения, с оплатой только за объём, с ИИ-ядром, "
           "SPV для международных инвесторов и выходом в ОАЭ."),
}
META = {
    "en": "Prepared by Consistente Ltd · Tallinn · info@consistente.tech · consistente.tech",
    "ru": "Подготовлено Consistente Ltd · Таллин · info@consistente.tech · consistente.tech",
}

# Each section is one slide. Types:
#   text    — eyebrow, title, bullets
#   table   — eyebrow, title, headers[], rows[][], (optional) note
#   diagram — eyebrow, title, diagram key, caption, (optional) bullets
#   shot    — eyebrow, title, img base (lang-prefixed), caption, (optional) bullets
SECTIONS = [
    # 1 — Executive summary
    {"type": "text",
     "eyebrow": {"en": "Executive summary", "ru": "Краткое резюме"},
     "title": {"en": "Three moves, one platform", "ru": "Три шага, одна платформа"},
     "bullets": {
        "en": [
            "**The gap.** Every Uzbek bank *rents* its factoring product from OzPlanet or FinMakon — Universalbank included (OzPlanet cooperation agreement, Sept 2024). None *owns* the product, the customer relationship or the data.",
            "**1 · Own the platform.** A white-label, AI-native platform Universalbank runs under its own brand — the bank owns the product, customers and data — priced on **financed volume only**.",
            "**2 · International capital.** An SPV module lets global — and specifically Gulf — investors fund Uzbek receivables. OzPlanet and FinMakon take **no foreign or retail money**, so this is a clean, uncontested differentiator.",
            "**3 · UAE entry.** A **DIFC SPV** and a phased Shariah programme (Murabaha → Wakala → Sukuk) to reach Dubai / DIFC Islamic capital.",
            "**Decomposable & live today.** Deploy origination first, add the investor / SPV layer later — and the chat-based loan triage and investor reporting are **already running** in the app.",
        ],
        "ru": [
            "**Разрыв.** Каждый узбекский банк *арендует* факторинговый продукт у OzPlanet или FinMakon — включая Universalbank (соглашение с OzPlanet, сент. 2024). Никто не *владеет* продуктом, отношениями с клиентом и данными.",
            "**1 · Владеть платформой.** White-label ИИ-платформа, которую Universalbank ведёт под своим брендом — банк владеет продуктом, клиентами и данными — с оплатой **только за профинансированный объём**.",
            "**2 · Международный капитал.** Модуль SPV позволяет глобальным — и прежде всего инвесторам Залива — финансировать узбекскую дебиторку. OzPlanet и FinMakon **не принимают иностранные или розничные деньги**, поэтому это чистый, незанятый дифференциатор.",
            "**3 · Выход в ОАЭ.** **SPV в DIFC** и поэтапная исламская программа (Мурабаха → Вакала → Сукук) для доступа к исламскому капиталу Дубая / DIFC.",
            "**Декомпозируемо и уже работает.** Сначала запускаем origination, слой инвесторов / SPV добавляем позже — а триаж заявок и отчётность в чате **уже работают** в приложении.",
        ],
     }},

    # 2 — The opportunity
    {"type": "text",
     "eyebrow": {"en": "The opportunity", "ru": "Возможность"},
     "title": {"en": "Uzbekistan built the perfect rails for invoice finance", "ru": "Узбекистан построил идеальную инфраструктуру для факторинга"},
     "bullets": {
        "en": [
            "**SoliqOnline** — the State Tax Committee's platform — has validated, timestamped and archived every B2B e-invoice in the country since **January 2020**.",
            "**Didox** connects **350,000+ organisations** to it and integrates with 1C — a ready-made supplier base that already exchanges invoices electronically.",
            "Every invoice is government-verified and buyer-confirmed on official record — **fraud and debtor-confirmation, factoring's hardest problems, are solved at source.**",
            "The market is scaling fast: the **CBU reports ~10.6 trillion soums (~$835M) factored in Jan–Sep 2025**, about **half of it digital** — against a **~$2bn** addressable market that is still lightly penetrated.",
        ],
        "ru": [
            "**SoliqOnline** — платформа Государственного налогового комитета — с **января 2020 года** проверяет, фиксирует по времени и архивирует каждую B2B-электронную счёт-фактуру в стране.",
            "**Didox** подключает к ней **350 000+ организаций** и интегрируется с 1С — готовая база поставщиков, уже обменивающихся счетами в электронном виде.",
            "Каждый счёт проверен государством и подтверждён покупателем в официальном реестре — **мошенничество и подтверждение дебитора, самые сложные проблемы факторинга, решены у источника.**",
            "Рынок быстро растёт: **ЦБ сообщает о ~10,6 трлн сум (~$835 млн) факторинга за янв–сен 2025**, около **половины — цифровой**, при адресном рынке **~$2 млрд** с всё ещё низким проникновением.",
        ],
     }},

    # 3 — Why now
    {"type": "text",
     "eyebrow": {"en": "Why now", "ru": "Почему сейчас"},
     "title": {"en": "The regulatory window is open", "ru": "Регуляторное окно открыто"},
     "bullets": {
        "en": [
            "**Presidential Decree No. 106 (Aug 2024)** mandated banks to offer factoring.",
            "**ZRU-1058 (Apr 2025)** amended the Civil Code to give factoring full legal force, and mandates automated API registration of the bank's collateral priority in the **CBU registry** — within about an hour of approval.",
            "**No new licence** is required: Universalbank's existing CBU licence covers this; Factorio operates as white-label technology under the bank's umbrella.",
            "**No Uzbek bank has yet built a production-grade SoliqOnline factoring stack with an AI core** — the first mover captures structural advantage.",
        ],
        "ru": [
            "**Указ Президента № 106 (авг. 2024)** обязал банки предоставлять факторинг.",
            "**ЗРУ-1058 (апр. 2025)** внёс поправки в Гражданский кодекс, придав факторингу полную юридическую силу, и обязал автоматически регистрировать залоговый приоритет банка в **реестре ЦБ РУз** по API — примерно в течение часа после одобрения.",
            "**Новая лицензия не требуется:** действующая лицензия ЦБ у Universalbank это покрывает; Factorio работает как white-label-технология под эгидой банка.",
            "**Ни один узбекский банк ещё не построил промышленный факторинговый стек на SoliqOnline с ИИ-ядром** — первопроходец получает структурное преимущество.",
        ],
     }},

    # 4 — The incumbents
    {"type": "table",
     "eyebrow": {"en": "The incumbents", "ru": "Действующие игроки"},
     "title": {"en": "OzPlanet & FinMakon own the market — as aggregators", "ru": "OzPlanet и FinMakon владеют рынком — как агрегаторы"},
     "headers": {
        "en": ["Dimension", "OzPlanet", "FinMakon", "Factorio (Consistente)"],
        "ru": ["Параметр", "OzPlanet", "FinMakon", "Factorio (Consistente)"],
     },
     "rows": {
        "en": [
            ["Model", "Aggregator; bank-set rates", "Aggregator (Didox-backed)", "**White-label marketplace** — bank owns it"],
            ["Digital share (Q3 2025)", "56% of digital volume", "44% of digital volume", "New entrant — builds share"],
            ["Funding sources", "Banks / MFOs only", "Banks / factoring cos", "**Banks + investors + retail + SPV**"],
            ["Foreign / retail capital", "**No**", "**No**", "**Yes** — SPV + marketplace"],
            ["Bank owns product & data", "No — OzPlanet does", "No — FinMakon does", "**Yes** — fully the bank's"],
            ["AI core", "Rule-based; nascent", "Nascent", "**Grok chat triage + reporting**, doc intelligence"],
            ["Pricing to bank", "Per-transaction commission", "Per-transaction commission", "**Volume-only bps** (no SaaS licence)"],
        ],
        "ru": [
            ["Модель", "Агрегатор; ставки банка", "Агрегатор (при Didox)", "**White-label маркетплейс** — владеет банк"],
            ["Доля в цифре (Q3 2025)", "56% цифрового объёма", "44% цифрового объёма", "Новичок — набирает долю"],
            ["Источники капитала", "Только банки / МФО", "Банки / факторинг-компании", "**Банки + инвесторы + розница + SPV**"],
            ["Иностр. / розн. капитал", "**Нет**", "**Нет**", "**Да** — SPV + маркетплейс"],
            ["Банк владеет продуктом и данными", "Нет — владеет OzPlanet", "Нет — владеет FinMakon", "**Да** — полностью банка"],
            ["ИИ-ядро", "Правила; зачаточно", "Зачаточно", "**Grok: триаж + отчётность**, обработка документов"],
            ["Цена для банка", "Комиссия за транзакцию", "Комиссия за транзакцию", "**Только б.п. от объёма** (без SaaS)"],
        ],
     },
     "note": {
        "en": "Universalbank already signed a cooperation agreement with OzPlanet (Sept 2024) — so it knows the model. The point of this proposal is not to rent a better aggregator, but to own the product outright. (Market data: CBU, Q3 2025.)",
        "ru": "Universalbank уже подписал соглашение о сотрудничестве с OzPlanet (сент. 2024) — то есть знает модель. Смысл этого предложения — не арендовать более удобный агрегатор, а владеть продуктом полностью. (Данные рынка: ЦБ, Q3 2025.)",
     }},

    # 5 — Own, don't rent (the key insight)
    {"type": "text",
     "eyebrow": {"en": "The strategic gap", "ru": "Стратегический разрыв"},
     "title": {"en": "No Uzbek bank owns its factoring product", "ru": "Ни один узбекский банк не владеет своим факторингом"},
     "bullets": {
        "en": [
            "Today a bank connecting to OzPlanet or FinMakon is a **funding partner on someone else's platform** — the aggregator owns the brand, the customer relationship and the data.",
            "That is a strategic dependency: the bank cannot differentiate, cannot cross-sell off its own data, and cannot switch without losing the clients.",
            "**Factorio flips this.** Universalbank runs the platform under its **own brand**, owns **all customer and transaction data**, and can export and move it — no platform lock-in.",
            "Same government rails (SoliqOnline, CBU registry), same speed — but the bank owns the asset instead of renting the rails. That ownership is the real product.",
        ],
        "ru": [
            "Сегодня банк, подключённый к OzPlanet или FinMakon, — это **партнёр по финансированию на чужой платформе**: агрегатор владеет брендом, отношениями с клиентом и данными.",
            "Это стратегическая зависимость: банк не может дифференцироваться, не может продавать по своим данным и не может сменить поставщика, не потеряв клиентов.",
            "**Factorio это переворачивает.** Universalbank ведёт платформу под **своим брендом**, владеет **всеми клиентскими и транзакционными данными** и может их выгрузить и перенести — без привязки к платформе.",
            "Та же гос. инфраструктура (SoliqOnline, реестр ЦБ), та же скорость — но банк владеет активом, а не арендует рельсы. Именно это владение и есть продукт.",
        ],
     }},

    # 6 — Money flow (bilingual diagram)
    {"type": "diagram", "bilingual": True,
     "eyebrow": {"en": "How it works", "ru": "Как это работает"},
     "title": {"en": "How the money moves — and who gets paid", "ru": "Как движутся деньги — и кто получает оплату"},
     "diagram": "15-money-flow",
     "caption": {"en": "Money & payment flow — seller, payer (debtor), platform, capital", "ru": "Поток денег и платежей — продавец, плательщик (дебитор), платформа, капитал"},
     "bullets": {
        "en": [
            "The **payer (debtor)** is the anchor: a specific, buyer-confirmed obligation on the SoliqOnline record.",
            "The platform advances ~85% now from the **capital source** (bank, investor or SPV); at maturity the payer settles 100% and the waterfall pays seller, capital and platform.",
            "Because the flow is identical regardless of who funds it, the **capital layer is pluggable** — bank first, investors/SPV later.",
        ],
        "ru": [
            "**Плательщик (дебитор)** — якорь: конкретное, подтверждённое покупателем обязательство в реестре SoliqOnline.",
            "Платформа выдаёт ~85% сразу из **источника капитала** (банк, инвестор или SPV); в срок плательщик гасит 100%, и водопад платит продавцу, капиталу и платформе.",
            "Поскольку поток одинаков независимо от источника финансирования, **слой капитала подключаемый** — сначала банк, инвесторы/SPV позже.",
        ],
     }},

    # 5 — Pillar 1 intro
    {"type": "text",
     "eyebrow": {"en": "Pillar 1 · A better platform", "ru": "Направление 1 · Лучшая платформа"},
     "title": {"en": "AI at the core, not bolted on", "ru": "ИИ в основе, а не сбоку"},
     "bullets": {
        "en": [
            "The same end-to-end factoring workflow — submit, verify, fund, settle — but with AI woven through triage, scoring, documents and reporting.",
            "Built on a lean, server-rendered stack (FastHTML + PostgreSQL) — fast to change, cheap to run, easy to white-label under the bank's brand.",
            "Trilingual by design (English · Oʻzbekcha · Russian); the AI answers in the user's language.",
            "The next three slides detail the AI, the credit-scoring engine, and the commercial model.",
        ],
        "ru": [
            "Тот же сквозной процесс факторинга — подача, проверка, финансирование, расчёт — но с ИИ в триаже, скоринге, документах и отчётности.",
            "На лёгком серверном стеке (FastHTML + PostgreSQL) — быстро менять, дёшево эксплуатировать, легко брендировать под банк.",
            "Трёхъязычность заложена в дизайн (английский · узбекский · русский); ИИ отвечает на языке пользователя.",
            "Следующие три слайда раскрывают ИИ, движок кредитного скоринга и коммерческую модель.",
        ],
     }},

    # Origination journey (bilingual diagram)
    {"type": "diagram", "bilingual": True,
     "eyebrow": {"en": "Journey · Borrower (origination)", "ru": "Путь · Заёмщик (origination)"},
     "title": {"en": "The seller's journey — sign-up to cash in 24–48h", "ru": "Путь продавца — от регистрации до денег за 24–48 ч"},
     "diagram": "16-origination-journey",
     "caption": {"en": "Origination: SoliqOnline import → AI triage → one-click approval → advance", "ru": "Origination: импорт из SoliqOnline → AI-триаж → одобрение в один клик → аванс"},
     "bullets": {
        "en": [
            "Manual bank involvement is a **single approval click** on a pre-populated screen; everything else is automated.",
            "The invoice is imported and verified from SoliqOnline; the debtor is scored; collateral is registered with the CBU in **under 60 minutes**.",
            "This is **Module A** — it stands alone and can go live first, funded entirely by the bank's balance sheet.",
        ],
        "ru": [
            "Ручное участие банка — **один клик одобрения** на заранее заполненном экране; всё остальное автоматизировано.",
            "Счёт импортируется и проверяется из SoliqOnline; дебитор скорится; залог регистрируется в ЦБ **менее чем за 60 минут**.",
            "Это **Модуль A** — он самодостаточен и может быть запущен первым, полностью за счёт баланса банка.",
        ],
     }},

    # 6 — AI chat flows (diagram)
    {"type": "diagram",
     "eyebrow": {"en": "Pillar 1 · AI", "ru": "Направление 1 · ИИ"},
     "title": {"en": "Two conversational surfaces, one Grok core", "ru": "Две разговорные поверхности, одно ядро Grok"},
     "diagram": "d2-ai-chat",
     "caption": {"en": "Chat-based loan triage and chat-based investor reporting", "ru": "Триаж заявок в чате и отчётность для инвесторов в чате"},
     "bullets": {
        "en": [
            "**Triage** turns a seller's plain-language description into an indicative grade, advance rate and document list — in seconds.",
            "**Reporting** answers an investor's questions grounded in their own live positions — no invented figures.",
            "Grok scores and explains; a human always approves. Every AI decision is logged and auditable.",
        ],
        "ru": [
            "**Триаж** превращает описание продавца на естественном языке в индикативный рейтинг, ставку аванса и список документов — за секунды.",
            "**Отчётность** отвечает на вопросы инвестора на основе его реальных позиций — без выдуманных цифр.",
            "Grok оценивает и объясняет; финальное решение всегда за человеком. Каждое решение ИИ логируется и проверяемо.",
        ],
     }},

    # 7 — Credit scoring (diagram)
    {"type": "diagram",
     "eyebrow": {"en": "Pillar 1 · Credit scoring", "ru": "Направление 1 · Кредитный скоринг"},
     "title": {"en": "Open-banking-style scoring, Uzbek edition", "ru": "Скоринг в духе open banking, узбекская версия"},
     "diagram": "d3-scoring",
     "caption": {"en": "Plaid-style data fusion adapted to Uzbekistan's rails", "ru": "Слияние данных в стиле Plaid, адаптированное к инфраструктуре Узбекистана"},
     "bullets": {
        "en": [
            "The Plaid model is *connect an account, read the cash flows, decide*. In Uzbekistan the richest feed is **SoliqOnline** — verified invoices and tax-declared turnover — plus **bank transaction data** and the **CRIF** bureau.",
            "A model produces the **grade, advance rate and price**; Grok writes the **plain-language rationale and adverse-action reasons**.",
            "Every decision is **versioned and reproducible** — Consistente's core methodology, and what a regulator will ask for.",
        ],
        "ru": [
            "Модель Plaid — это *подключи счёт, прочитай денежные потоки, прими решение*. В Узбекистане богатейший источник — **SoliqOnline** (проверенные счета и налоговый оборот) плюс **банковские транзакции** и бюро **CRIF**.",
            "Модель выдаёт **рейтинг, ставку аванса и цену**; Grok формирует **понятное обоснование и причины отказа**.",
            "Каждое решение **версионируется и воспроизводимо** — ключевая методология Consistente и то, что запросит регулятор.",
        ],
     }},

    # Economics of one invoice (table)
    {"type": "table",
     "eyebrow": {"en": "Economics", "ru": "Экономика"},
     "title": {"en": "What the bank earns on one invoice", "ru": "Что банк зарабатывает на одном счёте"},
     "headers": {
        "en": ["Cash flow · $10,000 invoice · 85% · 60d · 5%", "Amount", "Direction", "Day"],
        "ru": ["Поток · счёт $10 000 · 85% · 60 дн · 5%", "Сумма", "Направление", "День"],
     },
     "rows": {
        "en": [
            ["Advance to seller (85%)", "$8,500", "Platform → seller", "Day 1"],
            ["Payer (debtor) pays in full", "$10,000", "Payer → collection", "Day 60"],
            ["Net balance released to seller", "$1,500", "Platform → seller", "Day 60"],
            ["Discount income (gross)", "$500", "Capital keeps", "Day 60"],
            ["Platform fee (volume-based)", "≈ $13", "Bank → Consistente", "Day 60"],
            ["**Net income on the invoice**", "**≈ $487**", "**Capital keeps**", "**Day 60**"],
            ["**Annualised return on capital**", "**≈ 30% p.a.**", "—", "—"],
        ],
        "ru": [
            ["Аванс продавцу (85%)", "$8 500", "Платформа → продавец", "День 1"],
            ["Плательщик (дебитор) платит полностью", "$10 000", "Плательщик → сбор", "День 60"],
            ["Остаток продавцу", "$1 500", "Платформа → продавец", "День 60"],
            ["Дисконтный доход (брутто)", "$500", "Капитал удерживает", "День 60"],
            ["Комиссия платформы (за объём)", "≈ $13", "Банк → Consistente", "День 60"],
            ["**Чистый доход по счёту**", "**≈ $487**", "**Капитал удерживает**", "**День 60**"],
            ["**Годовая доходность на капитал**", "**≈ 30% год.**", "—", "—"],
        ],
     },
     "note": {
        "en": "Illustrative, per the standard SoliqOnline-verified structure. ~30% annualised on deployed capital compares with 22–28% on unsecured SME lending — at lower risk: the invoice is government-verified, the payer has confirmed the obligation, and the bank holds a registered CBU priority claim. Bank stays profitable to a ~4% default rate.",
        "ru": "Иллюстративно, по стандартной структуре с проверкой SoliqOnline. ~30% годовых на размещённый капитал против 22–28% по необеспеченному кредитованию МСБ — при меньшем риске: счёт проверен государством, плательщик подтвердил обязательство, банк держит зарегистрированный приоритет в ЦБ. Банк прибылен вплоть до ~4% дефолтов.",
     }},

    # 8 — Pricing (table)
    {"type": "table",
     "eyebrow": {"en": "Pillar 1 · Commercial model", "ru": "Направление 1 · Коммерческая модель"},
     "title": {"en": "Volume-only pricing — aligned with the bank", "ru": "Оплата только за объём — интересы совпадают с банком"},
     "headers": {
        "en": ["Monthly financed volume", "Platform fee (bps of financed volume)"],
        "ru": ["Месячный объём финансирования", "Комиссия платформы (б.п. от объёма)"],
     },
     "rows": {
        "en": [
            ["Up to UZS 50 bn", "120 bps"],
            ["UZS 50–200 bn", "90 bps"],
            ["UZS 200–500 bn", "70 bps"],
            ["Over UZS 500 bn", "55 bps"],
        ],
        "ru": [
            ["До 50 млрд сум", "120 б.п."],
            ["50–200 млрд сум", "90 б.п."],
            ["200–500 млрд сум", "70 б.п."],
            ["Свыше 500 млрд сум", "55 б.п."],
        ],
     },
     "note": {
        "en": "Illustrative. No SaaS licence, no per-seat, no setup fee — Consistente is paid only when the bank finances an invoice. Options on the table: a revenue-share alternative (a small % of the bank's discount income instead of bps), 12-month bank exclusivity in Uzbekistan, and a Year-3 deferred-equity option structured as new shares (non-dilutive, no board/veto rights). International SPV module: ~60 bps p.a. on invested AUM + a performance share above a hurdle. Final figures to be set with Universalbank.",
        "ru": "Иллюстративно. Без SaaS-лицензии, без оплаты за место, без платы за внедрение — Consistente получает вознаграждение только когда банк финансирует счёт. Опции на столе: альтернатива с revenue-share (небольшой % дисконтного дохода банка вместо б.п.), 12-месячная эксклюзивность для банка в Узбекистане и отложенный опцион на долю в Год 3 в виде новых акций (без размытия, без прав в совете/вето). Модуль международного SPV: ~60 б.п. в год от инвестированных активов + доля от результата сверх порога. Итоговые цифры — по согласованию с Universalbank.",
     }},

    # 9 — Pillar 2 intro (diagram)
    {"type": "diagram",
     "eyebrow": {"en": "Pillar 2 · International capital", "ru": "Направление 2 · Международный капитал"},
     "title": {"en": "An SPV that opens Uzbek receivables to the world", "ru": "SPV, открывающее узбекскую дебиторку миру"},
     "diagram": "d4-spv",
     "caption": {"en": "Invoice pool → Factorio → DIFC SPV → international investors", "ru": "Пул счетов → Factorio → SPV в DIFC → международные инвесторы"},
     "bullets": {
        "en": [
            "Universalbank remains **originator and servicer**; a bankruptcy-remote **SPV** holds the investor-facing interest and channels foreign capital into the invoice pool.",
            "Two tracks off the **same asset base**: a **conventional** note/participation for institutional investors, and **Shariah** structures for Gulf capital.",
            "Investors get the same grounded, AI-assisted reporting — in their language — plus statements and a clean audit trail.",
        ],
        "ru": [
            "Universalbank остаётся **оригинатором и сервисером**; защищённое от банкротства **SPV** держит инвесторскую долю и направляет иностранный капитал в пул счетов.",
            "Два трека на **одной базе активов**: **конвенциональная** нота/участие для институционалов и **исламские** структуры для капитала Залива.",
            "Инвесторы получают ту же обоснованную ИИ-отчётность — на своём языке — плюс выписки и прозрачный аудиторский след.",
        ],
     }},

    # Investor journey (bilingual diagram)
    {"type": "diagram", "bilingual": True,
     "eyebrow": {"en": "Journey · Investor", "ru": "Путь · Инвестор"},
     "title": {"en": "The investor's journey — onboarding to distribution", "ru": "Путь инвестора — от онбординга до распределения"},
     "diagram": "17-investor-journey",
     "caption": {"en": "Investor: onboard → commit capital → hold with AI reporting → get paid", "ru": "Инвестор: онбординг → капитал → удержание с AI-отчётностью → выплата"},
     "bullets": {
        "en": [
            "Retail, institutional and Gulf investors onboard once (KYC / accreditation), then fund via the **marketplace** or the **SPV**.",
            "They hold positions with **on-demand AI portfolio reporting**; at maturity, principal + return are distributed automatically.",
            "This is **Module B** — it plugs onto the same origination engine later, without changing anything a seller or the bank does.",
        ],
        "ru": [
            "Розничные, институциональные и инвесторы Залива проходят онбординг один раз (KYC / аккредитация), затем финансируют через **маркетплейс** или **SPV**.",
            "Они держат позиции с **AI-отчётностью по запросу**; в срок основная сумма и доход распределяются автоматически.",
            "Это **Модуль B** — он подключается к тому же движку origination позже, ничего не меняя для продавца или банка.",
        ],
     }},

    # Decomposition (bilingual diagram)
    {"type": "diagram", "bilingual": True,
     "eyebrow": {"en": "Delivery · Decomposability", "ru": "Реализация · Декомпозируемость"},
     "title": {"en": "Two modules — buy one, or both, in sequence", "ru": "Два модуля — берите один или оба, последовательно"},
     "diagram": "18-decomposition",
     "caption": {"en": "Origination (Module A) and Investors/SPV (Module B) are independently deployable", "ru": "Origination (Модуль A) и Инвесторы/SPV (Модуль B) внедряются независимо"},
     "bullets": {
        "en": [
            "**Module A — Origination** can be Universalbank's whole first phase: bank-funded factoring, live in weeks, immediately satisfying the Decree-106 mandate.",
            "**Module B — Investors & SPV** adds the marketplace and DIFC/international capital later, over the same engine, with no rework.",
            "The **capital layer is the seam**: swap or add funding sources without touching origination — de-risking the programme and the investment.",
        ],
        "ru": [
            "**Модуль A — Origination** может стать всей первой фазой Universalbank: факторинг за счёт банка, запуск за недели, немедленное выполнение требований Указа № 106.",
            "**Модуль B — Инвесторы и SPV** добавляет маркетплейс и капитал DIFC/международный позже, поверх того же движка, без переделок.",
            "**Слой капитала — это стык**: меняйте или добавляйте источники финансирования, не трогая origination — снижая риск программы и инвестиций.",
        ],
     }},

    # 10 — Pillar 3: instruments (table)
    {"type": "table",
     "eyebrow": {"en": "Pillar 3 · UAE entry", "ru": "Направление 3 · Выход в ОАЭ"},
     "title": {"en": "A phased Shariah programme for DIFC capital", "ru": "Поэтапная исламская программа для капитала DIFC"},
     "headers": {
        "en": ["Criterion", "Murabaha SCF", "Wakala Fund", "Sukuk", "Musharaka"],
        "ru": ["Критерий", "Мурабаха (SCF)", "Фонд Вакала", "Сукук", "Мушарака"],
     },
     "rows": {
        "en": [
            ["Shariah purity", "High", "High", "High", "Highest"],
            ["Complexity", "Low", "Low–Med", "High", "Medium"],
            ["Time to market", "Fastest", "Fast", "Slowest", "Medium"],
            ["Dubai investor appeal", "Medium", "High", "Very high", "High"],
            ["Capital per deal", "Small–Med", "Medium", "Large", "Med–Large"],
            ["Regulatory need", "Fatwa only", "Fatwa + fund", "Fatwa + DFSA", "Fatwa + fund"],
            ["Best for", "Pilot", "Family offices", "Institutional", "Strategic partners"],
        ],
        "ru": [
            ["Шариатская чистота", "Высокая", "Высокая", "Высокая", "Наивысшая"],
            ["Сложность", "Низкая", "Низк.–ср.", "Высокая", "Средняя"],
            ["Срок вывода", "Самый быстрый", "Быстро", "Дольше всех", "Средне"],
            ["Привлекательность (Дубай)", "Средняя", "Высокая", "Очень высокая", "Высокая"],
            ["Капитал на сделку", "Мал.–ср.", "Средний", "Крупный", "Ср.–крупн."],
            ["Регуляторика", "Только фатва", "Фатва + фонд", "Фатва + DFSA", "Фатва + фонд"],
            ["Кому подходит", "Пилот", "Family office", "Институционалы", "Стратег. партнёры"],
        ],
     },
     "note": {
        "en": "Recommended path: run them in sequence — Murabaha SCF pilot (M6–12) → Wakala fund → Sukuk. Each phase builds the track record that makes the next credible. All require an AAOIFI-accredited fatwa before deployment.",
        "ru": "Рекомендуемый путь: последовательно — пилот Мурабаха (M6–12) → фонд Вакала → Сукук. Каждый этап формирует историю, делающую следующий убедительным. Все требуют фатвы аккредитованного AAOIFI совета до запуска.",
     }},

    # 11 — Why Dubai
    {"type": "text",
     "eyebrow": {"en": "Pillar 3 · Why Dubai / DIFC", "ru": "Направление 3 · Почему Дубай / DIFC"},
     "title": {"en": "The yield spread is the story", "ru": "Главное — спред доходности"},
     "bullets": {
        "en": [
            "Uzbek factoring yields of **~20–30% p.a.** against Gulf Islamic money-market returns of **~4–6%** — an exceptional spread for the risk, given SoliqOnline's structural protections.",
            "**Short duration** (30–90-day rolling receivables) is rare and highly demanded by Gulf liquidity managers.",
            "Uzbek receivables are **uncorrelated** with Gulf real estate, regional equities or oil — genuine diversification.",
            "**DIFC (DFSA)** and **ADGM (FSRA)** both have mature Islamic-finance frameworks and recognise SPV/Sukuk structures; the UAE's Federal Decree-Law No. 50 of 2022 codifies the contracts.",
        ],
        "ru": [
            "Доходность узбекского факторинга **~20–30% годовых** против исламского денежного рынка Залива **~4–6%** — исключительный спред для такого риска благодаря структурной защите SoliqOnline.",
            "**Короткая дюрация** (оборот 30–90 дней) редка и очень востребована у ликвидити-менеджеров Залива.",
            "Узбекская дебиторка **не коррелирует** с недвижимостью Залива, региональными акциями или нефтью — настоящая диверсификация.",
            "**DIFC (DFSA)** и **ADGM (FSRA)** имеют зрелые рамки исламских финансов и признают структуры SPV/Сукук; Федеральный декрет-закон № 50 от 2022 кодифицирует контракты.",
        ],
     }},

    # 12 — Architecture overview (diagram)
    {"type": "diagram",
     "eyebrow": {"en": "Architecture", "ru": "Архитектура"},
     "title": {"en": "One process, AI-native, integration-ready", "ru": "Один процесс, ИИ-нативный, готовый к интеграциям"},
     "diagram": "d1-architecture",
     "caption": {"en": "Target system architecture with AI components highlighted", "ru": "Целевая архитектура системы с выделенными ИИ-компонентами"},
     "bullets": {
        "en": [
            "A single FastHTML process serves the landing site, the investor app and the AI assistants; PostgreSQL holds the factoring and AI data.",
            "Clean integration seams to **SoliqOnline/Didox**, **open banking**, **CRIF**, the **CBU registry**, and the **DIFC SPV / custodian**.",
            "Grok (x.ai) is reached over an OpenAI-compatible interface — the model id is configuration, avoiding vendor lock-in.",
        ],
        "ru": [
            "Единый процесс FastHTML обслуживает лендинг, приложение инвестора и ИИ-ассистентов; PostgreSQL хранит данные факторинга и ИИ.",
            "Чистые точки интеграции с **SoliqOnline/Didox**, **open banking**, **CRIF**, **реестром ЦБ** и **SPV/кастодианом в DIFC**.",
            "Grok (x.ai) подключается через OpenAI-совместимый интерфейс — id модели задаётся конфигурацией, что исключает привязку к вендору.",
        ],
     }},

    # 13 — Proof: prototype (shot)
    {"type": "shot",
     "eyebrow": {"en": "Proof · Working today", "ru": "Доказательство · Работает уже сегодня"},
     "title": {"en": "The AI is not a slide — it's shipped", "ru": "ИИ — не слайд, он уже внедрён"},
     "img": "13-ai-triage",
     "caption": {"en": "Live chat-based invoice triage in the Factorio app", "ru": "Действующий триаж счёта в чате в приложении Factorio"},
     "bullets": {
        "en": [
            "A seller describes an invoice in a sentence; the assistant returns an indicative risk band, advance rate and next-document list.",
            "Built on Grok (x.ai), trilingual, with graceful fallback and full audit logging — the same pattern extends to scoring and reporting.",
        ],
        "ru": [
            "Продавец описывает счёт одним предложением; ассистент возвращает индикативный рейтинг, ставку аванса и список документов.",
            "На базе Grok (x.ai), трёхъязычно, с корректным запасным поведением и полным аудитом — тот же подход расширяется на скоринг и отчётность.",
        ],
     }},

    # 14 — Roadmap (diagram)
    {"type": "diagram",
     "eyebrow": {"en": "Delivery", "ru": "Реализация"},
     "title": {"en": "From prototype to regional platform", "ru": "От прототипа к региональной платформе"},
     "diagram": "d5-roadmap",
     "caption": {"en": "Phased rollout — each phase funds the next", "ru": "Поэтапный запуск — каждый этап финансирует следующий"},
     "bullets": {
        "en": [
            "**Phase 0 (now):** AI triage + reporting prototype live.",
            "**Phases 1–2 (M1–6):** platform GA, SoliqOnline integration, open-banking scoring, volume pricing.",
            "**Phases 3–4 (M6+):** DIFC SPV and first international capital, then the Wakala → Sukuk sequence and regional scale.",
        ],
        "ru": [
            "**Этап 0 (сейчас):** прототип ИИ-триажа и отчётности работает.",
            "**Этапы 1–2 (M1–6):** релиз платформы, интеграция SoliqOnline, скоринг на open banking, оплата за объём.",
            "**Этапы 3–4 (M6+):** SPV в DIFC и первый международный капитал, затем последовательность Вакала → Сукук и региональный масштаб.",
        ],
     }},

    # 15 — Why Consistente
    {"type": "text",
     "eyebrow": {"en": "Why Consistente", "ru": "Почему Consistente"},
     "title": {"en": "Production AI, delivered consistently", "ru": "Промышленный ИИ, поставляемый стабильно"},
     "bullets": {
        "en": [
            "Consistente Ltd (Tallinn, EU) builds **production-grade AI for enterprises** — with reproducible pipelines, versioned models and inspectable prompts, not black boxes.",
            "Precedents across **financial services and regulated sectors**: LSEG, DBRS Morningstar, ARM, Microsoft.",
            "Core capabilities map directly onto this project: **document intelligence**, **applied forecasting/scoring**, and **agentic workflows** with human review.",
            "EU-based, audit-first — the right profile for a bank building AI a regulator will scrutinise.",
        ],
        "ru": [
            "Consistente Ltd (Таллин, ЕС) создаёт **промышленный ИИ для предприятий** — с воспроизводимыми пайплайнами, версионируемыми моделями и прозрачными промптами, а не «чёрные ящики».",
            "Опыт в **финансовых услугах и регулируемых секторах**: LSEG, DBRS Morningstar, ARM, Microsoft.",
            "Ключевые компетенции прямо ложатся на проект: **интеллектуальная обработка документов**, **прикладной прогноз/скоринг** и **агентные процессы** с проверкой человеком.",
            "Базируется в ЕС, приоритет аудируемости — верный профиль для банка, строящего ИИ под надзором регулятора.",
        ],
     }},

    # 16 — Next steps
    {"type": "text",
     "eyebrow": {"en": "Next steps", "ru": "Следующие шаги"},
     "title": {"en": "What we propose to do first", "ru": "Что предлагаем сделать первым"},
     "bullets": {
        "en": [
            "**1.** Agree the white-label scope and the volume-based commercial terms.",
            "**2.** Stand up a pilot on live SoliqOnline data; ship the platform GA with the AI triage and scoring engine.",
            "**3.** In parallel, begin DIFC SPV and fatwa preparation so international capital can follow the operating book.",
            "**Contact:** Consistente Ltd · info@consistente.tech · consistente.tech",
        ],
        "ru": [
            "**1.** Согласовать объём white-label и коммерческие условия на основе объёма.",
            "**2.** Запустить пилот на реальных данных SoliqOnline; вывести платформу в релиз с ИИ-триажем и движком скоринга.",
            "**3.** Параллельно начать подготовку SPV в DIFC и фатвы, чтобы международный капитал последовал за операционным портфелем.",
            "**Контакт:** Consistente Ltd · info@consistente.tech · consistente.tech",
        ],
     }},
]


# ── Markdown ────────────────────────────────────────────────────────────────

def _bullets_md(items):
    return "\n".join(f"- {b}" for b in items)


def _table_md(headers, rows):
    out = ["| " + " | ".join(headers) + " |",
           "| " + " | ".join("---" for _ in headers) + " |"]
    for r in rows:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _slide_md(s, lang, *, for_html):
    """Render one section to markdown. for_html swaps mermaid fences for images."""
    parts = [f'<p class="eyebrow">{s["eyebrow"][lang]}</p>', "", f'## {s["title"][lang]}', ""]
    t = s["type"]
    if t == "text":
        parts += [_bullets_md(s["bullets"][lang])]
    elif t == "table":
        parts += [_table_md(s["headers"][lang], s["rows"][lang])]
        if s.get("note"):
            parts += ["", f'> {s["note"][lang]}']
    elif t == "diagram":
        key = s["diagram"]
        bilingual = s.get("bilingual")
        img_rel = f'img/{lang}-{key}.png' if bilingual else f'img/proposal/{key}.png'
        src = BILINGUAL[key][lang] if bilingual else DIAGRAMS[key]
        if for_html:
            parts += [f'![{s["caption"][lang]}]({img_rel})',
                      "", f'<p class="caption">{s["caption"][lang]}</p>']
        else:
            parts += ["```mermaid", src.strip(), "```",
                      "", f'*{s["caption"][lang]}*']
        if s.get("bullets"):
            parts += ["", _bullets_md(s["bullets"][lang])]
    elif t == "shot":
        parts += [f'![{s["caption"][lang]} screen](img/{lang}-{s["img"]}.png)',
                  "", f'<p class="caption">{s["caption"][lang]}</p>']
        if s.get("bullets"):
            parts += ["", _bullets_md(s["bullets"][lang])]
    return "\n".join(parts)


def build_markdown(lang, *, for_html=False):
    parts = [f"# {TITLE[lang]}", f"## {SUBTITLE[lang]}", "", META[lang], "", INTRO[lang], "", "---", ""]
    for s in SECTIONS:
        parts += [_slide_md(s, lang, for_html=for_html), "", "---", ""]
    return "\n".join(parts).rstrip("- \n") + "\n"


def build_html(lang):
    body = md_lib.markdown(build_markdown(lang, for_html=True),
                           extensions=["tables", "attr_list", "sane_lists", "fenced_code"])
    title = f"{TITLE[lang]}"
    return (f'<!doctype html>\n<html lang="{lang}">\n<head>\n<meta charset="utf-8">\n'
            f'<title>{title}</title>\n<link rel="stylesheet" href="assets/proposal.css">\n'
            f'</head>\n<body>\n{body}\n</body>\n</html>\n')


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _strip_md(text):
    return text.replace("**", "").replace("`", "").replace("*", "")


def _add_bg(slide, prs):
    shp = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shp.fill.solid(); shp.fill.fore_color.rgb = PARCH
    shp.line.fill.background(); shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)


def _title_block(slide, eyebrow, title):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12), Inches(0.4))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = eyebrow.upper()
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = GOLD
    bar = slide.shapes.add_shape(1, Inches(0.7), Inches(0.8), Inches(11.9), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background(); bar.shadow.inherit = False
    btf = bar.text_frame; btf.word_wrap = True; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    br = btf.paragraphs[0].add_run(); br.text = "   " + title
    br.font.size = Pt(21); br.font.bold = True; br.font.color.rgb = WHITE


def _bullets_box(slide, bullets, left, top, width, height, size=15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        run = p.add_run(); run.text = "•  " + _strip_md(b)
        run.font.size = Pt(size); run.font.color.rgb = INK


def _fit_picture(slide, path, box_l, box_t, box_w, box_h):
    """Add a picture scaled to fit within a box, centred."""
    from PIL import Image
    iw, ih = Image.open(path).size
    box_ar, img_ar = box_w / box_h, iw / ih
    if img_ar > box_ar:
        w = box_w; h = int(box_w / img_ar)
    else:
        h = box_h; w = int(box_h * img_ar)
    left = box_l + (box_w - w) // 2
    top = box_t + (box_h - h) // 2
    slide.shapes.add_picture(str(path), left, top, width=w, height=h)


def _add_table(slide, headers, rows, left, top, width):
    nrows, ncols = len(rows) + 1, len(headers)
    height = Inches(0.4) * nrows
    gtbl = slide.shapes.add_table(nrows, ncols, left, top, width, height).table
    for c, h in enumerate(headers):
        cell = gtbl.cell(0, c); cell.text = _strip_md(h)
        para = cell.text_frame.paragraphs[0]; para.runs[0].font.size = Pt(11)
        para.runs[0].font.bold = True; para.runs[0].font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = GREEN
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtbl.cell(r, c); cell.text = _strip_md(val)
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].font.size = Pt(10.5); para.runs[0].font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else RGBColor(0xEE, 0xED, 0xE5)


def build_pptx(lang, out):
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank); _add_bg(s, prs)
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.4))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = TITLE[lang]
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = GREEN
    p = tf.add_paragraph(); r = p.add_run(); r.text = SUBTITLE[lang]
    r.font.size = Pt(20); r.font.color.rgb = MUTED
    p = tf.add_paragraph(); r = p.add_run(); r.text = _strip_md(INTRO[lang])
    r.font.size = Pt(13); r.font.color.rgb = INK
    p = tf.add_paragraph(); r = p.add_run(); r.text = _strip_md(META[lang])
    r.font.size = Pt(11); r.font.color.rgb = MUTED

    for sec in SECTIONS:
        s = prs.slides.add_slide(blank); _add_bg(s, prs)
        _title_block(s, sec["eyebrow"][lang], sec["title"][lang])
        t = sec["type"]
        if t == "text":
            _bullets_box(s, sec["bullets"][lang], Inches(0.75), Inches(2.0), Inches(11.9), Inches(5.0))
        elif t == "table":
            _add_table(s, sec["headers"][lang], sec["rows"][lang], Inches(0.75), Inches(2.0), Inches(11.85))
            if sec.get("note"):
                nb = s.shapes.add_textbox(Inches(0.75), Inches(6.6), Inches(11.9), Inches(0.8))
                nb.text_frame.word_wrap = True
                run = nb.text_frame.paragraphs[0].add_run(); run.text = _strip_md(sec["note"][lang])
                run.font.size = Pt(10); run.font.italic = True; run.font.color.rgb = MUTED
        elif t in ("diagram", "shot"):
            has_bul = bool(sec.get("bullets"))
            if t == "diagram":
                if sec.get("bilingual"):
                    img_path = IMG / f"{lang}-{sec['diagram']}.png"
                else:
                    img_path = DIAG / f"{sec['diagram']}.png"
            else:
                img_path = IMG / f"{lang}-{sec['img']}.png"
            img_w = Inches(7.2) if has_bul else Inches(12.4)
            if img_path.exists():
                _fit_picture(s, img_path, Inches(0.6), Inches(2.0), img_w, Inches(4.9))
            cap = s.shapes.add_textbox(Inches(0.6), Inches(6.95), img_w, Inches(0.4))
            cr = cap.text_frame.paragraphs[0].add_run(); cr.text = _strip_md(sec["caption"][lang])
            cr.font.size = Pt(10); cr.font.italic = True; cr.font.color.rgb = MUTED
            cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            if has_bul:
                _bullets_box(s, sec["bullets"][lang], Inches(8.1), Inches(2.0), Inches(4.9), Inches(5.0), size=13)

    prs.save(str(out))


# ── Driver ───────────────────────────────────────────────────────────────────

def build_lang(lang):
    base = DOCS / f"universalbank_consistente_proposal_{lang}"
    (base.with_suffix(".md")).write_text(build_markdown(lang), encoding="utf-8")
    html = build_html(lang)
    (base.with_suffix(".html")).write_text(html, encoding="utf-8")
    HTML(string=html, base_url=str(DOCS)).write_pdf(str(base.with_suffix(".pdf")))
    build_pptx(lang, base.with_suffix(".pptx"))
    for ext in ("md", "html", "pdf", "pptx"):
        f = base.with_suffix("." + ext)
        print(f"  {f.name}  ({f.stat().st_size/1024:.0f} KB)")


def main():
    print("Rendering diagrams…")
    render_diagrams()
    render_bilingual()
    print("Building Universalbank proposal")
    for lang in ("en", "ru"):
        print(f"→ {lang}")
        build_lang(lang)
    print("done")


if __name__ == "__main__":
    main()
