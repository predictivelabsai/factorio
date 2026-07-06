"""Build the Factorio user guide in four formats, per language (en / ru).

For each language it writes, into docs/, a date-slugged set:
    factorio_user_guide_<YYYY-MM-DD>_<lang>.md     (source slide deck)
    factorio_user_guide_<YYYY-MM-DD>_<lang>.html   (markdown + assets/guide.css)
    factorio_user_guide_<YYYY-MM-DD>_<lang>.pdf     (WeasyPrint, A4 landscape slides)
    factorio_user_guide_<YYYY-MM-DD>_<lang>.pptx    (python-pptx, 16:9 deck)

One slide per screen; screenshots live in docs/img/<lang>-NN-*.png. Content is
authored once below (bilingual) so all formats and both languages stay in sync.

    python -m scripts.build_user_guide
"""

from __future__ import annotations

from datetime import date
from pathlib import Path

import markdown as md_lib
from weasyprint import HTML
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
IMG = DOCS / "img"
CSS = DOCS / "assets" / "guide.css"
DATE = date.today().isoformat()
SITE = "factorio.co.uk"

# Brand
GREEN = RGBColor(0x1F, 0x5D, 0x43)
INK = RGBColor(0x14, 0x23, 0x1B)
MUTED = RGBColor(0x41, 0x50, 0x46)
GOLD = RGBColor(0xC8, 0xA2, 0x4B)
PARCH = RGBColor(0xF7, 0xF6, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Content (authored once, bilingual) ───────────────────────────────────

TITLE = {
    "en": "Factorio — User Guide",
    "ru": "Factorio — Руководство пользователя",
}
SUBTITLE = {
    "en": "Invoice financing, end to end",
    "ru": "Факторинг счетов — от начала до конца",
}
INTRO = {
    "en": ("Factorio is an invoice-financing (factoring) marketplace. Businesses sell their "
           "unpaid invoices for cash today; investors fund those invoices for a short-term, "
           "asset-backed return. This guide walks through every screen — the public site and "
           "the investor product app — one screen per slide."),
    "ru": ("Factorio — это маркетплейс факторинга (финансирования счетов). Компании продают "
           "свои неоплаченные счета и получают деньги сегодня; инвесторы финансируют эти счета "
           "ради краткосрочного дохода, обеспеченного активами. Это руководство показывает "
           "каждый экран — публичный сайт и инвесторское приложение — по одному экрану на слайд."),
}
META = {
    "en": f"Generated {DATE} · Live at **{SITE}**",
    "ru": f"Сформировано {DATE} · Доступно на **{SITE}**",
}

# Each slide: img base (lang prefix added at build), eyebrow, title, bullets, caption.
SLIDES = [
    {
        "img": "01-landing",
        "eyebrow": {"en": "Overview", "ru": "Обзор"},
        "title": {"en": "What Factorio is", "ru": "Что такое Factorio"},
        "caption": {"en": "Factorio homepage", "ru": "Главная страница Factorio"},
        "bullets": {
            "en": [
                "An invoice-financing marketplace connecting two sides: **sellers** who need cash now and **investors** seeking short-term, asset-backed returns.",
                "One site serves both — a marketing front and an HTMX-driven product app.",
                "The hero summarizes the model: advance rate, days-to-funding, sectors served and total funded volume (UZS).",
                "Fully trilingual — English, Oʻzbekcha and Russian — switched from the top navigation.",
            ],
            "ru": [
                "Маркетплейс факторинга, соединяющий две стороны: **продавцов**, которым нужны деньги сейчас, и **инвесторов**, ищущих краткосрочный доход, обеспеченный активами.",
                "Один сайт обслуживает обе стороны — маркетинговую витрину и приложение на HTMX.",
                "Первый экран отражает модель: ставку аванса, срок финансирования, отрасли и общий объём финансирования (UZS).",
                "Полностью трёхъязычный — английский, узбекский и русский — переключается в верхней навигации.",
            ],
        },
    },
    {
        "img": "02-for-sellers",
        "eyebrow": {"en": "For sellers", "ru": "Для продавцов"},
        "title": {"en": "Turn invoices into working capital", "ru": "Превратите счета в оборотный капитал"},
        "caption": {"en": "For Sellers page", "ru": "Страница «Для продавцов»"},
        "bullets": {
            "en": [
                "Get a large share of an invoice advanced in days instead of waiting 30–120 days for the debtor to pay.",
                "The seller journey: submit an invoice → get it verified and risk-graded → receive the advance → settle when the debtor pays.",
                "Factoring, not a loan — no new debt on the balance sheet; funding scales with sales.",
                "Built for manufacturing, wholesale, construction, logistics and services.",
            ],
            "ru": [
                "Получите значительную часть суммы счёта авансом за несколько дней вместо ожидания оплаты 30–120 дней.",
                "Путь продавца: подать счёт → проверка и присвоение рейтинга риска → получение аванса → расчёт после оплаты должником.",
                "Это факторинг, а не кредит — без новой задолженности; финансирование растёт вместе с продажами.",
                "Создано для производства, оптовой торговли, строительства, логистики и услуг.",
            ],
        },
    },
    {
        "img": "03-for-investors",
        "eyebrow": {"en": "For investors", "ru": "Для инвесторов"},
        "title": {"en": "Short-term, asset-backed returns", "ru": "Краткосрочный доход, обеспеченный активами"},
        "caption": {"en": "For Investors page", "ru": "Страница «Для инвесторов»"},
        "bullets": {
            "en": [
                "Fund verified invoices and earn a fee/return over a short hold — often just weeks.",
                "Diversify across debtors, sectors and risk grades (A–D).",
                "Returns are tied to real receivables, not market speculation.",
                "Leads into the investor app — dashboard, marketplace and portfolio.",
            ],
            "ru": [
                "Финансируйте проверенные счета и получайте доход за короткий срок — часто всего несколько недель.",
                "Диверсификация по должникам, отраслям и рейтингам риска (A–D).",
                "Доход привязан к реальной дебиторской задолженности, а не к рыночным спекуляциям.",
                "Ведёт в инвесторское приложение — дашборд, маркетплейс и портфель.",
            ],
        },
    },
    {
        "img": "04-how-it-works",
        "eyebrow": {"en": "How it works", "ru": "Как это работает"},
        "title": {"en": "Four steps, end to end", "ru": "Четыре шага, от начала до конца"},
        "caption": {"en": "How It Works page", "ru": "Страница «Как это работает»"},
        "bullets": {
            "en": [
                "**Submit** — a seller uploads an invoice and company details.",
                "**Verify** — Factorio checks the debtor and assigns a risk grade and pricing.",
                "**Fund** — the invoice appears on the marketplace; investors fund it to its goal.",
                "**Settle** — when the debtor pays, principal + interest are distributed and the position closes.",
            ],
            "ru": [
                "**Подача** — продавец загружает счёт и данные компании.",
                "**Проверка** — Factorio проверяет должника и присваивает рейтинг риска и цену.",
                "**Финансирование** — счёт появляется на маркетплейсе; инвесторы финансируют его до цели.",
                "**Расчёт** — после оплаты должником распределяются основная сумма и проценты, позиция закрывается.",
            ],
        },
    },
    {
        "img": "15-money-flow",
        "is_diagram": True,
        "eyebrow": {"en": "How the money moves", "ru": "Как движутся деньги"},
        "title": {"en": "Where every party's money comes from", "ru": "Откуда каждая сторона получает деньги"},
        "caption": {"en": "Money & payment flow — seller, payer, platform, capital",
                    "ru": "Поток денег и платежей — продавец, плательщик, платформа, капитал"},
        "bullets": {
            "en": [
                "The **seller** delivers goods and issues a SoliqOnline e-invoice; the **payer (debtor)** owes the invoice amount.",
                "The platform advances ~85% to the seller now, funded by the **capital source** — the bank's balance sheet, investors, or the DIFC SPV.",
                "At maturity the **payer pays 100%** into the collection account; the seller receives the net balance, the capital source gets principal + return, and the platform keeps a volume-based fee.",
                "The flow is identical whoever funds it — which is what makes the investor / SPV layer pluggable.",
            ],
            "ru": [
                "**Продавец** поставляет товар и выставляет э-счёт в SoliqOnline; **плательщик (дебитор)** должен сумму счёта.",
                "Платформа выдаёт продавцу ~85% сразу, за счёт **источника капитала** — баланса банка, инвесторов или SPV в DIFC.",
                "В срок **плательщик платит 100%** на счёт сбора; продавец получает остаток, источник капитала — основную сумму и доход, платформа удерживает комиссию за объём.",
                "Поток одинаков независимо от того, кто финансирует — именно поэтому слой инвесторов / SPV подключаемый.",
            ],
        },
    },
    {
        "img": "16-origination-journey",
        "is_diagram": True,
        "eyebrow": {"en": "Journey · Borrower (origination)", "ru": "Путь · Заёмщик (origination)"},
        "title": {"en": "The seller's journey, end to end", "ru": "Путь продавца, от начала до конца"},
        "caption": {"en": "From sign-up to advance in 24–48 hours", "ru": "От регистрации до аванса за 24–48 часов"},
        "bullets": {
            "en": [
                "Sign up (bank KYC) → invoice auto-imported from SoliqOnline → request financing in an AI-triage chat.",
                "Debtor scoring returns indicative terms; the bank approves in one click; collateral is registered with the CBU in under an hour.",
                "The advance lands in the seller's account in 24–48 hours — no branch visit, no paper.",
            ],
            "ru": [
                "Регистрация (KYC банка) → счёт импортируется из SoliqOnline → запрос финансирования в чате AI-триажа.",
                "Скоринг дебитора возвращает индикативные условия; банк одобряет в один клик; залог регистрируется в ЦБ менее чем за час.",
                "Аванс поступает на счёт продавца за 24–48 часов — без визита в отделение и без бумаг.",
            ],
        },
    },
    {
        "img": "17-investor-journey",
        "is_diagram": True,
        "eyebrow": {"en": "Journey · Investor", "ru": "Путь · Инвестор"},
        "title": {"en": "The investor's journey, end to end", "ru": "Путь инвестора, от начала до конца"},
        "caption": {"en": "From onboarding to distribution", "ru": "От онбординга до распределения дохода"},
        "bullets": {
            "en": [
                "Onboard (KYC / accreditation) → commit capital via the marketplace or the SPV → auto-invest or pick invoices.",
                "Hold the position with AI portfolio reporting on demand; when the debtor pays at maturity, principal + return are distributed.",
                "**Decomposable:** the origination piece (bank-funded) can launch first; the investor marketplace and SPV plug in later, independently.",
            ],
            "ru": [
                "Онбординг (KYC / аккредитация) → внесение капитала через маркетплейс или SPV → автоинвест или выбор счетов.",
                "Позиция удерживается с AI-отчётностью по запросу; когда дебитор платит в срок, распределяются основная сумма и доход.",
                "**Декомпозируемо:** часть origination (за счёт банка) можно запустить первой; маркетплейс инвесторов и SPV подключаются позже, независимо.",
            ],
        },
    },
    {
        "img": "05-pricing",
        "eyebrow": {"en": "Pricing", "ru": "Тарифы"},
        "title": {"en": "Transparent, per-invoice fees", "ru": "Прозрачные тарифы за каждый счёт"},
        "caption": {"en": "Pricing page", "ru": "Страница «Тарифы»"},
        "bullets": {
            "en": [
                "Pricing is per invoice — a fee per 30 days against the advanced amount, with no hidden charges.",
                "Shows the headline advance rate and worked example economics.",
                "No subscription and no commitment to submit an invoice.",
                "A clear breakdown so both sides understand the return and the cost.",
            ],
            "ru": [
                "Тариф рассчитывается за каждый счёт — комиссия за 30 дней от суммы аванса, без скрытых платежей.",
                "Показаны базовая ставка аванса и наглядный пример расчёта.",
                "Без подписки и без обязательств при подаче счёта.",
                "Понятная разбивка, чтобы обе стороны видели доход и стоимость.",
            ],
        },
    },
    {
        "img": "06-contact",
        "eyebrow": {"en": "Contact", "ru": "Контакты"},
        "title": {"en": "Get in touch", "ru": "Связаться с нами"},
        "caption": {"en": "Contact page", "ru": "Страница «Контакты»"},
        "bullets": {
            "en": [
                "A contact page for sellers and investors to reach the Factorio team.",
                "Shows the contact email and a simple enquiry form.",
                "The entry point to start onboarding.",
            ],
            "ru": [
                "Страница контактов для продавцов и инвесторов, чтобы связаться с командой Factorio.",
                "Показаны контактный e-mail и простая форма обращения.",
                "Точка входа для начала подключения.",
            ],
        },
    },
    {
        "img": "07-dashboard",
        "eyebrow": {"en": "Investor app · Dashboard", "ru": "Приложение · Дашборд"},
        "title": {"en": "Your personalized overview", "ru": "Ваш персональный обзор"},
        "caption": {"en": "Investor dashboard", "ru": "Дашборд инвестора"},
        "bullets": {
            "en": [
                "The `/app` home greets the current investor and leads with their own KPIs: portfolio value, net annual return, earned-to-date and the next settlement.",
                "A recent-activity feed surfaces notifications — fundings, settlements and updates.",
                "Quick actions jump to the marketplace, portfolio and statement.",
                "Platform-wide stats are kept as a secondary strip.",
                "Top-right **investor switcher** lets you view the app as any investor (demo, no password).",
            ],
            "ru": [
                "Главная `/app` приветствует текущего инвестора и показывает его KPI: стоимость портфеля, чистую годовую доходность, заработанное и ближайший расчёт.",
                "Лента активности отражает уведомления — финансирования, расчёты и обновления.",
                "Быстрые действия ведут в маркетплейс, портфель и выписку.",
                "Общая статистика платформы вынесена во вторичную полосу.",
                "**Переключатель инвестора** справа вверху позволяет смотреть приложение от лица любого инвестора (демо, без пароля).",
            ],
        },
    },
    {
        "img": "08-marketplace",
        "eyebrow": {"en": "Investor app · Marketplace", "ru": "Приложение · Маркетплейс"},
        "title": {"en": "Browse fundable invoices", "ru": "Просмотр счетов для финансирования"},
        "caption": {"en": "Marketplace listings", "ru": "Список на маркетплейсе"},
        "bullets": {
            "en": [
                "Every open, fundable invoice is a card: debtor, sector, risk grade, amount and a live funding-progress bar.",
                "Each card shows the key economics — advance rate, fee per 30 days and estimated return — plus the target hold in days.",
                "The filter bar narrows by sector, risk grade, maximum term and minimum return.",
                "Click a card to open the full invoice detail.",
            ],
            "ru": [
                "Каждый открытый счёт — это карточка: должник, отрасль, рейтинг риска, сумма и шкала прогресса финансирования.",
                "На карточке видны ключевые показатели — ставка аванса, комиссия за 30 дней и ожидаемая доходность — а также срок удержания в днях.",
                "Панель фильтров сужает выбор по отрасли, рейтингу риска, максимальному сроку и минимальной доходности.",
                "Нажмите на карточку, чтобы открыть подробности счёта.",
            ],
        },
    },
    {
        "img": "09-marketplace-detail",
        "eyebrow": {"en": "Investor app · Invoice detail", "ru": "Приложение · Детали счёта"},
        "title": {"en": "Full transparency before you invest", "ru": "Полная прозрачность перед инвестицией"},
        "caption": {"en": "Invoice detail page", "ru": "Страница деталей счёта"},
        "bullets": {
            "en": [
                "The headline shows the debtor, risk grade, sector and invoice number.",
                "A funding panel: amount raised vs goal, advance rate, fee and estimated return.",
                "A **debtor-company profile** lists registration number, sector, country and annual turnover.",
                "Everything needed to make a funding decision on a single screen.",
            ],
            "ru": [
                "В заголовке — должник, рейтинг риска, отрасль и номер счёта.",
                "Панель финансирования: собрано против цели, ставка аванса, комиссия и ожидаемая доходность.",
                "Блок **профиля компании-должника**: регистрационный номер, отрасль, страна и годовой оборот.",
                "Всё необходимое для решения о финансировании на одном экране.",
            ],
        },
    },
    {
        "img": "10-portfolio",
        "eyebrow": {"en": "Investor app · Portfolio", "ru": "Приложение · Портфель"},
        "title": {"en": "The reporting cockpit", "ru": "Аналитическая панель"},
        "caption": {"en": "Portfolio cockpit", "ru": "Панель портфеля"},
        "bullets": {
            "en": [
                "Two hero panels: **net annual return** (interest received, write-offs, net income) and **account value** (active invested, expected outstanding, realized).",
                "An **investments aging table** buckets active positions by days to / past due.",
                "A **payment-habits table** shows how settled invoices actually paid — early, on time, late or written off.",
                "A positions table lists every investment with realized return, settlement date, grade and status.",
            ],
            "ru": [
                "Две главные панели: **чистая годовая доходность** (полученные проценты, списания, чистый доход) и **стоимость счёта** (активные вложения, ожидаемое к получению, реализованное).",
                "**Таблица старения вложений** распределяет активные позиции по сроку до / после погашения.",
                "**Таблица платёжной дисциплины** показывает, как фактически платили по закрытым счетам — раньше, в срок, с опозданием или списание.",
                "Таблица позиций перечисляет каждую инвестицию с реализованным доходом, датой расчёта, рейтингом и статусом.",
            ],
        },
    },
    {
        "img": "11-statement",
        "eyebrow": {"en": "Investor app · Account statement", "ru": "Приложение · Выписка по счёту"},
        "title": {"en": "Every transaction, filterable", "ru": "Каждая операция, с фильтрами"},
        "caption": {"en": "Account statement", "ru": "Выписка по счёту"},
        "bullets": {
            "en": [
                "A unified ledger of cash movements: investments out (−) and settlements in (+).",
                "Filter by date range, type, counterparty, invoice number and minimum amount.",
                "Signed amounts in UZS, newest first.",
                "One-click **CSV export** for bookkeeping and reconciliation.",
            ],
            "ru": [
                "Единый реестр движения средств: вложения (−) и поступления по расчётам (+).",
                "Фильтрация по периоду, типу, контрагенту, номеру счёта и минимальной сумме.",
                "Суммы со знаком в UZS, новые сверху.",
                "**Экспорт в CSV** в один клик для учёта и сверки.",
            ],
        },
    },
    {
        "img": "12-auto-invest",
        "eyebrow": {"en": "Investor app · Auto-invest", "ru": "Приложение · Авто-инвестирование"},
        "title": {"en": "Set rules, invest automatically", "ru": "Задайте правила — инвестируйте автоматически"},
        "caption": {"en": "Auto-invest settings", "ru": "Настройки авто-инвестирования"},
        "bullets": {
            "en": [
                "Configure automated bidding: minimum risk grade, maximum amount per invoice and preferred sectors.",
                "Toggle the strategy active/inactive; the status is shown at the top.",
                "Rules are saved per investor and applied to matching new invoices.",
                "Modelled on investly.co's autobidder.",
            ],
            "ru": [
                "Настройте автоматические заявки: минимальный рейтинг риска, максимальную сумму на счёт и предпочитаемые отрасли.",
                "Включайте и выключайте стратегию; статус показан сверху.",
                "Правила сохраняются для каждого инвестора и применяются к подходящим новым счетам.",
                "Создано по образцу автобиддера investly.co.",
            ],
        },
    },
    {
        "img": "13-ai-triage",
        "eyebrow": {"en": "AI · Application triage", "ru": "AI · Триаж заявок"},
        "title": {"en": "Chat-based invoice triage", "ru": "Триаж счёта в чате"},
        "caption": {"en": "AI triage assistant", "ru": "AI-ассистент триажа"},
        "bullets": {
            "en": [
                "Sellers describe an invoice and debtor in plain language; the assistant collects only what is still missing.",
                "It returns an **indicative risk band (A–D)**, an **advance rate** and the **documents the bank needs next** — in seconds, not days.",
                "Powered by **Grok (x.ai)**; everything is indicative and subject to verification — no final decision is automated.",
                "Fully trilingual — the assistant replies in the visitor's language (English, Oʻzbekcha, Russian).",
            ],
            "ru": [
                "Продавцы описывают счёт и должника обычными словами; ассистент запрашивает только недостающее.",
                "Он возвращает **индикативный уровень риска (A–D)**, **ставку аванса** и **список документов для банка** — за секунды, а не за дни.",
                "На базе **Grok (x.ai)**; все данные индикативны и подлежат проверке — окончательное решение не автоматизируется.",
                "Полностью трёхъязычный — ассистент отвечает на языке посетителя (английский, узбекский, русский).",
            ],
        },
    },
    {
        "img": "14-ai-assistant",
        "eyebrow": {"en": "AI · Portfolio reporting", "ru": "AI · Отчёты по портфелю"},
        "title": {"en": "Ask your portfolio anything", "ru": "Спросите свой портфель о чём угодно"},
        "caption": {"en": "AI reporting assistant", "ru": "AI-ассистент отчётности"},
        "bullets": {
            "en": [
                "Investors ask questions in natural language — net return, exposure by sector or debtor, overdue positions, upcoming settlements.",
                "Answers are **grounded in the investor's own live positions** — every figure traces back to real portfolio data, not invented.",
                "Concrete and specific: it quotes actual amounts and invoice numbers, and says so plainly when the data can't answer.",
                "The conversational layer on top of the portfolio cockpit — reporting on demand, in the investor's language.",
            ],
            "ru": [
                "Инвесторы задают вопросы на естественном языке — чистая доходность, доля по отраслям или должникам, просроченные позиции, ближайшие расчёты.",
                "Ответы **основаны на реальных позициях инвестора** — каждая цифра восходит к данным портфеля, а не выдумана.",
                "Конкретно и точно: приводятся реальные суммы и номера счетов, а при нехватке данных об этом честно сообщается.",
                "Разговорный слой поверх аналитической панели портфеля — отчётность по запросу, на языке инвестора.",
            ],
        },
    },
]

CLOSING = {
    "eyebrow": {"en": "Get started", "ru": "Начало работы"},
    "title": {"en": "Start with Factorio", "ru": "Начните с Factorio"},
    "bullets": {
        "en": [
            "**Sellers** — submit an invoice and get an offer in 1–3 working days.",
            "**Investors** — browse the marketplace, build a portfolio, automate with auto-invest.",
            f"Live at **{SITE}** · hello@factorio.co.uk",
        ],
        "ru": [
            "**Продавцам** — подайте счёт и получите предложение за 1–3 рабочих дня.",
            "**Инвесторам** — изучайте маркетплейс, формируйте портфель, автоматизируйте через авто-инвестирование.",
            f"Доступно на **{SITE}** · hello@factorio.co.uk",
        ],
    },
}


# ── Markdown ──────────────────────────────────────────────────────────────

def _bullets_md(items: list[str]) -> str:
    return "\n".join(f"- {b}" for b in items)


def build_markdown(lang: str) -> str:
    parts = [
        f"# {TITLE[lang]}",
        f"## {SUBTITLE[lang]}",
        "",
        META[lang],
        "",
        INTRO[lang],
        "",
        "---",
        "",
    ]
    for s in SLIDES:
        if s.get("is_diagram"):
            # full-width, centred flow/journey diagram (not a floated screenshot)
            img_md = (f'![{s["caption"][lang]}](img/{lang}-{s["img"]}.png){{: .diagram }}\n\n'
                      f'<p class="caption">{s["caption"][lang]}</p>')
        else:
            img_md = f'![{s["caption"][lang]}](img/{lang}-{s["img"]}.png)'
        parts += [
            f'<p class="eyebrow">{s["eyebrow"][lang]}</p>',
            "",
            f'## {s["title"][lang]}',
            "",
            img_md,
            "",
            _bullets_md(s["bullets"][lang]),
            "",
            "---",
            "",
        ]
    parts += [
        f'<p class="eyebrow">{CLOSING["eyebrow"][lang]}</p>',
        "",
        f'## {CLOSING["title"][lang]}',
        "",
        _bullets_md(CLOSING["bullets"][lang]),
        "",
    ]
    return "\n".join(parts)


def build_html(md_text: str, lang: str) -> str:
    body = md_lib.markdown(md_text, extensions=["tables", "attr_list", "sane_lists"])
    title = f"{TITLE[lang]} ({DATE})"
    return (
        f'<!doctype html>\n<html lang="{lang}">\n<head>\n'
        f'<meta charset="utf-8">\n<title>{title}</title>\n'
        f'<link rel="stylesheet" href="assets/guide.css">\n'
        f'</head>\n<body>\n{body}\n</body>\n</html>\n'
    )


# ── PPTX ───────────────────────────────────────────────────────────────────

def _strip_md(text: str) -> str:
    return text.replace("**", "").replace("`", "")


def _add_bg(slide, prs):
    shp = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shp.fill.solid(); shp.fill.fore_color.rgb = PARCH
    shp.line.fill.background()
    shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


def _textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    return tb, tf


def build_pptx(lang: str, out: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # title slide
    s = prs.slides.add_slide(blank)
    _add_bg(s, prs)
    _, tf = _textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3))
    p = tf.paragraphs[0]; r = p.add_run(); r.text = TITLE[lang]
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = GREEN
    p2 = tf.add_paragraph(); r = p2.add_run(); r.text = SUBTITLE[lang]
    r.font.size = Pt(24); r.font.color.rgb = MUTED
    p3 = tf.add_paragraph(); r = p3.add_run()
    r.text = _strip_md(META[lang]); r.font.size = Pt(13); r.font.color.rgb = MUTED

    # content slides
    for sl in SLIDES + [{"img": None, **CLOSING}]:
        s = prs.slides.add_slide(blank)
        _add_bg(s, prs)
        # eyebrow
        _, tf = _textbox(s, Inches(0.7), Inches(0.4), Inches(12), Inches(0.4))
        r = tf.paragraphs[0].add_run(); r.text = sl["eyebrow"][lang].upper()
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = GOLD
        # title bar
        bar = s.shapes.add_shape(1, Inches(0.7), Inches(0.8), Inches(11.9), Inches(0.85))
        bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
        bar.shadow.inherit = False
        btf = bar.text_frame; btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run(); br.text = "   " + sl["title"][lang]
        br.font.size = Pt(22); br.font.bold = True; br.font.color.rgb = WHITE

        has_img = sl["img"] is not None
        is_diagram = sl.get("is_diagram", False)
        img_path = (IMG / f"{lang}-{sl['img']}.png") if has_img else None

        if is_diagram:
            # wide flow/journey diagram centred on top, bullets across the bottom
            if img_path and img_path.exists():
                s.shapes.add_picture(str(img_path), Inches(1.67), Inches(1.9), width=Inches(10.0))
            _, tf = _textbox(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(2.2))
            for i, b in enumerate(sl["bullets"][lang]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run(); run.text = "•  " + _strip_md(b)
                run.font.size = Pt(13); run.font.color.rgb = INK
        else:
            text_w = Inches(5.6) if has_img else Inches(11.9)
            _, tf = _textbox(s, Inches(0.7), Inches(2.0), text_w, Inches(5.0))
            for i, b in enumerate(sl["bullets"][lang]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                run = p.add_run(); run.text = "•  " + _strip_md(b)
                run.font.size = Pt(15); run.font.color.rgb = INK
            if img_path and img_path.exists():
                s.shapes.add_picture(str(img_path), Inches(6.5), Inches(2.0), width=Inches(6.3))

    prs.save(str(out))


# ── Driver ─────────────────────────────────────────────────────────────────

def build_lang(lang: str) -> None:
    base = DOCS / f"factorio_user_guide_{DATE}_{lang}"
    md_text = build_markdown(lang)
    (base.with_suffix(".md")).write_text(md_text, encoding="utf-8")

    html = build_html(md_text, lang)
    (base.with_suffix(".html")).write_text(html, encoding="utf-8")

    # base_url = docs/ so img/ and assets/ resolve
    HTML(string=html, base_url=str(DOCS)).write_pdf(str(base.with_suffix(".pdf")))

    build_pptx(lang, base.with_suffix(".pptx"))

    for ext in ("md", "html", "pdf", "pptx"):
        f = base.with_suffix("." + ext)
        print(f"  {f.name}  ({f.stat().st_size/1024:.0f} KB)")


def main() -> None:
    print(f"Building Factorio user guide ({DATE})")
    print("Rendering flow/journey diagrams…")
    from scripts.diagrams import render_bilingual
    render_bilingual()
    for lang in ("en", "ru"):
        print(f"→ {lang}")
        build_lang(lang)
    print("done")


if __name__ == "__main__":
    main()
